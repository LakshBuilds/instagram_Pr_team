#!/usr/bin/env python3
"""
Script to convert Markdown files to PowerPoint and Word documents
Requires: python-pptx, python-docx, markdown
"""

import os
import sys
from pathlib import Path

def install_requirements():
    """Install required packages"""
    packages = ['python-pptx', 'python-docx', 'markdown']
    
    for package in packages:
        try:
            __import__(package.replace('-', '_'))
            print(f"✅ {package} already installed")
        except ImportError:
            print(f"📦 Installing {package}...")
            os.system(f"pip install {package}")

def convert_markdown_to_pptx(md_file, output_file):
    """Convert markdown presentation outline to PowerPoint"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # Create presentation
        prs = Presentation()
        
        # Read markdown file
        with open(md_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Split by slides (## Slide markers)
        slides_content = content.split('## Slide')[1:]  # Skip header
        
        for i, slide_content in enumerate(slides_content):
            lines = slide_content.strip().split('\n')
            
            # Extract slide title
            title_line = lines[0] if lines else f"Slide {i+1}"
            title = title_line.split(':')[0].strip()
            
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            slide.shapes.title.text = title
            
            # Add content
            if len(lines) > 1:
                content_text = '\n'.join(lines[1:])
                # Clean up markdown formatting
                content_text = content_text.replace('**', '').replace('*', '').replace('```', '')
                content_text = content_text.replace('- ', '• ')
                
                # Add to content placeholder
                if slide.shapes.placeholders[1].has_text_frame:
                    slide.shapes.placeholders[1].text = content_text
        
        # Save presentation
        prs.save(output_file)
        print(f"✅ PowerPoint created: {output_file}")
        
    except Exception as e:
        print(f"❌ Error creating PowerPoint: {e}")

def convert_markdown_to_docx(md_file, output_file):
    """Convert markdown report to Word document"""
    try:
        from docx import Document
        from docx.shared import Inches
        
        # Create document
        doc = Document()
        
        # Read markdown file
        with open(md_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Split by sections
        sections = content.split('\n## ')
        
        # Add title
        title = doc.add_heading('Instagram Reels Analytics Dashboard - Project Report', 0)
        
        for section in sections[1:]:  # Skip first empty section
            lines = section.strip().split('\n')
            
            # Add section heading
            if lines:
                heading = lines[0].strip()
                doc.add_heading(heading, level=1)
                
                # Add content
                current_para = ""
                for line in lines[1:]:
                    line = line.strip()
                    
                    if line.startswith('### '):
                        # Sub-heading
                        if current_para:
                            doc.add_paragraph(current_para)
                            current_para = ""
                        doc.add_heading(line[4:], level=2)
                    elif line.startswith('**') and line.endswith('**'):
                        # Bold text
                        if current_para:
                            doc.add_paragraph(current_para)
                            current_para = ""
                        p = doc.add_paragraph()
                        p.add_run(line[2:-2]).bold = True
                    elif line.startswith('- '):
                        # Bullet point
                        if current_para:
                            doc.add_paragraph(current_para)
                            current_para = ""
                        doc.add_paragraph(line[2:], style='List Bullet')
                    elif line.startswith('```'):
                        # Code block
                        continue
                    elif line:
                        # Regular text
                        current_para += line + " "
                
                if current_para:
                    doc.add_paragraph(current_para)
        
        # Save document
        doc.save(output_file)
        print(f"✅ Word document created: {output_file}")
        
    except Exception as e:
        print(f"❌ Error creating Word document: {e}")

def main():
    print("🚀 Converting Markdown files to Office formats...")
    
    # Install requirements
    install_requirements()
    
    # File paths
    presentation_md = "PRESENTATION_OUTLINE.md"
    report_md = "PROJECT_REPORT.md"
    
    output_pptx = "Instagram_Reels_Analytics_Presentation.pptx"
    output_docx = "Instagram_Reels_Analytics_Report.docx"
    
    # Check if markdown files exist
    if not os.path.exists(presentation_md):
        print(f"❌ {presentation_md} not found!")
        return
    
    if not os.path.exists(report_md):
        print(f"❌ {report_md} not found!")
        return
    
    # Convert files
    print(f"📄 Converting {presentation_md} to PowerPoint...")
    convert_markdown_to_pptx(presentation_md, output_pptx)
    
    print(f"📄 Converting {report_md} to Word document...")
    convert_markdown_to_docx(report_md, output_docx)
    
    print("\n🎉 Conversion completed!")
    print(f"📁 Files created:")
    print(f"   • {output_pptx}")
    print(f"   • {output_docx}")
    print("\n💡 You can now open these files and copy content to your original format files.")

if __name__ == "__main__":
    main()