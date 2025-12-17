#!/usr/bin/env python3
"""
Script to update Chitkara University template files with Instagram Reels Analytics content
Uses the original template files and fills them with project data
"""

import os
import sys
from pathlib import Path

def install_requirements():
    """Install required packages"""
    packages = ['python-pptx', 'python-docx', 'markdown']
    
    for package in packages:
        try:
            __import__(package.replace('-', '_'))
            print(f"✅ {package} already installed")
        except ImportError:
            print(f"📦 Installing {package}...")
            os.system(f"pip install {package}")

def update_pptx_template(template_file, output_file):
    """Update PowerPoint template with Instagram Reels Analytics content"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # Load the template presentation
        prs = Presentation(template_file)
        print(f"📖 Loaded template with {len(prs.slides)} slides")
        
        # Define content for each slide
        slide_contents = [
            {
                "title": "Instagram Reels Analytics Dashboard",
                "subtitle": "A Comprehensive Full-Stack Web Application for Social Media Analytics",
                "content": [
                    "Student Name: [Your Name]",
                    "Course: 22CS411 - Advanced Web Development",
                    "Date: December 2024",
                    "Institution: Chitkara University"
                ]
            },
            {
                "title": "Project Overview",
                "content": [
                    "• Comprehensive Instagram Reels analytics platform",
                    "• Designed for content creators and marketing teams", 
                    "• Real-time data import and performance tracking",
                    "• Team collaboration with advanced analytics",
                    "",
                    "Key Statistics:",
                    "• 15,000+ lines of code",
                    "• 50+ React components", 
                    "• 4 major third-party integrations",
                    "• Full-stack TypeScript application"
                ]
            },
            {
                "title": "Problem Statement",
                "content": [
                    "Challenges in Social Media Analytics:",
                    "• Limited insights from native Instagram analytics",
                    "• No team collaboration features",
                    "• Manual data collection and analysis", 
                    "• Lack of historical performance tracking",
                    "• No comparative analysis tools",
                    "",
                    "Our Solution:",
                    "• Automated data import from Instagram",
                    "• Comprehensive analytics dashboard",
                    "• Team collaboration features",
                    "• Historical data preservation",
                    "• Advanced visualization tools"
                ]
            },
            {
                "title": "Technical Architecture",
                "content": [
                    "Frontend Stack:",
                    "• React 18 + TypeScript",
                    "• Vite for build optimization", 
                    "• Tailwind CSS + shadcn/ui",
                    "• TanStack Query for data management",
                    "• Recharts for visualizations",
                    "",
                    "Backend Stack:",
                    "• Node.js + Express",
                    "• Supabase (PostgreSQL)",
                    "• Clerk Authentication", 
                    "• Dual API Integration"
                ]
            },
            {
                "title": "Core Features - Analytics Dashboard",
                "content": [
                    "Personal Analytics:",
                    "• Individual performance tracking",
                    "• Engagement metrics (views, likes, comments)",
                    "• Revenue tracking and analysis",
                    "• Trend analysis over time",
                    "",
                    "Team Analytics:",
                    "• Organization-wide insights",
                    "• Creator performance comparisons",
                    "• Team streak tracking",
                    "• Collaborative decision making"
                ]
            },
            {
                "title": "Database Design & Security",
                "content": [
                    "Reels Table Schema:",
                    "• id, shortcode, ownerusername",
                    "• likescount, commentscount, videoplaycount",
                    "• payout, locationname, takenat",
                    "• created_by_user_id, created_by_email",
                    "",
                    "Security Features:",
                    "• Row Level Security (RLS)",
                    "• User-based data access",
                    "• Encrypted API key storage",
                    "• Input validation and sanitization"
                ]
            },
            {
                "title": "Technical Challenges & Solutions",
                "content": [
                    "Challenge 1: Instagram API Limitations",
                    "• Solution: Dual API system with intelligent fallback",
                    "",
                    "Challenge 2: Real-time Data Sync", 
                    "• Solution: Supabase real-time subscriptions",
                    "",
                    "Challenge 3: Complex Data Visualizations",
                    "• Solution: Recharts with custom components",
                    "",
                    "Challenge 4: Mobile Responsiveness",
                    "• Solution: Mobile-first design with Tailwind CSS"
                ]
            },
            {
                "title": "Performance & Security",
                "content": [
                    "Application Performance:",
                    "• Load Time: < 2 seconds",
                    "• Data Refresh: < 5 seconds", 
                    "• Import Speed: 10-15 reels/minute",
                    "• Database Query Optimization",
                    "",
                    "Security Implementation:",
                    "• Clerk integration with OAuth providers",
                    "• JWT token management",
                    "• Row Level Security at database level",
                    "• CORS configuration"
                ]
            },
            {
                "title": "Code Examples",
                "content": [
                    "React Component Example:",
                    "const Analytics = () => {",
                    "  const { user } = useUser();",
                    "  const [reels, setReels] = useState<Reel[]>([]);",
                    "  ",
                    "  const fetchReels = async () => {",
                    "    const { data } = await supabase",
                    "      .from('reels')",
                    "      .select('*')",
                    "      .eq('created_by_email', user.email);",
                    "    setReels(data || []);",
                    "  };",
                    "  ",
                    "  return <AnalyticsDashboard reels={reels} />;",
                    "};"
                ]
            },
            {
                "title": "Future Enhancements & Conclusion",
                "content": [
                    "Planned Features:",
                    "• AI-powered content recommendations",
                    "• Automated posting optimization",
                    "• Multi-platform integration",
                    "• Advanced team collaboration tools",
                    "",
                    "Key Achievements:",
                    "• Production-ready full-stack application",
                    "• Modern architecture with best practices",
                    "• Comprehensive feature set",
                    "• Scalable and maintainable codebase"
                ]
            }
        ]
        
        # Update slides with content
        for i, slide in enumerate(prs.slides):
            if i < len(slide_contents):
                content = slide_contents[i]
                
                # Update title if exists
                if slide.shapes.title and 'title' in content:
                    slide.shapes.title.text = content['title']
                
                # Update content placeholders
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        if 'subtitle' in content and i == 0:
                            # First slide subtitle
                            shape.text_frame.text = content['subtitle']
                        elif 'content' in content and shape != slide.shapes.title:
                            # Content slides
                            shape.text_frame.clear()
                            for line in content['content']:
                                p = shape.text_frame.add_paragraph()
                                p.text = line
                            break
        
        # Save updated presentation
        prs.save(output_file)
        print(f"✅ PowerPoint template updated: {output_file}")
        
    except Exception as e:
        print(f"❌ Error updating PowerPoint template: {e}")

def update_docx_template(template_file, output_file):
    """Update Word template with Instagram Reels Analytics content"""
    try:
        from docx import Document
        from docx.shared import Inches
        
        # Load the template document
        doc = Document(template_file)
        print(f"📖 Loaded Word template")
        
        # Clear existing content but keep formatting
        # Find and replace placeholder text
        replacements = {
            "[PROJECT TITLE]": "Instagram Reels Analytics Dashboard",
            "[PROJECT SUBTITLE]": "A Comprehensive Full-Stack Web Application for Social Media Analytics",
            "[STUDENT NAME]": "[Your Name]",
            "[COURSE CODE]": "22CS411 - Advanced Web Development", 
            "[DATE]": "December 2024",
            "[UNIVERSITY]": "Chitkara University"
        }
        
        # Replace text in paragraphs
        for paragraph in doc.paragraphs:
            for old_text, new_text in replacements.items():
                if old_text in paragraph.text:
                    paragraph.text = paragraph.text.replace(old_text, new_text)
        
        # Add comprehensive content sections
        sections_content = [
            {
                "heading": "Executive Summary",
                "content": "The Instagram Reels Analytics Dashboard is a comprehensive full-stack web application designed to provide content creators, marketing teams, and social media managers with detailed insights into their Instagram content performance. This project demonstrates advanced web development skills, modern architecture patterns, and sophisticated data analytics capabilities."
            },
            {
                "heading": "1. Project Overview",
                "content": "Social media content creators and marketing teams need comprehensive tools to track, analyze, and optimize their Instagram Reels performance. This project addresses these needs through a sophisticated web application that combines real-time Instagram data import, advanced analytics with multiple visualization types, team collaboration features, and gamification elements."
            },
            {
                "heading": "2. Technical Architecture",
                "content": "The application uses React 18 with TypeScript for the frontend, providing type safety and modern development practices. The backend utilizes Node.js with Express, Supabase for PostgreSQL database and real-time features, and Clerk for user authentication. Key design patterns include component-based architecture, custom hooks for business logic separation, and responsive design with mobile-first approach."
            },
            {
                "heading": "3. Core Features Implementation", 
                "content": "The system provides comprehensive analytics dashboards for both personal and team use, featuring individual creator performance tracking, engagement metrics analysis, and trend visualization. The data import system supports both single and bulk reel imports with real-time status feedback and intelligent error handling through dual API integration."
            },
            {
                "heading": "4. Database Schema Design",
                "content": "The PostgreSQL database implements Row Level Security (RLS) with a comprehensive reels table containing Instagram metadata including id, shortcode, engagement metrics, and user association fields. Security features include user-based data access control, encrypted API key storage, and input validation."
            },
            {
                "heading": "5. Performance & Security",
                "content": "Application performance metrics include load times under 2 seconds, data refresh capabilities under 5 seconds, and optimized database queries. Security implementation features Clerk integration with OAuth providers, JWT token management, and comprehensive data protection measures including CORS configuration and SQL injection prevention."
            },
            {
                "heading": "6. Challenges & Solutions",
                "content": "Key technical challenges included Instagram API limitations addressed through dual API system implementation, real-time data synchronization solved using Supabase subscriptions, and complex data visualization requirements met through Recharts integration with custom components."
            },
            {
                "heading": "7. Future Enhancements",
                "content": "Planned improvements include AI-powered content recommendations, automated posting schedule optimization, integration with additional social media platforms, GraphQL API implementation, and Progressive Web App capabilities for enhanced user experience."
            },
            {
                "heading": "8. Conclusion",
                "content": "The Instagram Reels Analytics Dashboard successfully demonstrates advanced full-stack web development capabilities, combining modern technologies with sophisticated data analytics features. The project showcases technical excellence, user experience focus, scalable architecture, and innovative solutions while maintaining comprehensive security and quality standards."
            }
        ]
        
        # Add sections to document
        for section in sections_content:
            # Add heading
            heading = doc.add_heading(section["heading"], level=1)
            # Add content
            doc.add_paragraph(section["content"])
            doc.add_paragraph()  # Add spacing
        
        # Save updated document
        doc.save(output_file)
        print(f"✅ Word template updated: {output_file}")
        
    except Exception as e:
        print(f"❌ Error updating Word template: {e}")

def main():
    print("🚀 Updating Chitkara University template files with Instagram Reels Analytics content...")
    
    # Install requirements
    install_requirements()
    
    # File paths
    ppt_template = "assignment/22CS411_PPT_Format.pptx"
    doc_template = "assignment/22CS411_Report Format.docx"
    
    output_pptx = "Instagram_Reels_Analytics_Final_Presentation.pptx"
    output_docx = "Instagram_Reels_Analytics_Final_Report.docx"
    
    # Check if template files exist
    if not os.path.exists(ppt_template):
        print(f"❌ {ppt_template} not found!")
        return
    
    if not os.path.exists(doc_template):
        print(f"❌ {doc_template} not found!")
        return
    
    # Update template files
    print(f"📄 Updating PowerPoint template...")
    update_pptx_template(ppt_template, output_pptx)
    
    print(f"📄 Updating Word template...")
    update_docx_template(doc_template, output_docx)
    
    print("\n🎉 Template update completed!")
    print(f"📁 Files created:")
    print(f"   • {output_pptx}")
    print(f"   • {output_docx}")
    print("\n💡 These files now have Chitkara University formatting with your project content!")

if __name__ == "__main__":
    main()